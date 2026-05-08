from __future__ import annotations

import csv
from pathlib import Path
from typing import Any

import pandas as pd
from docx import Document
from pptx import Presentation
from pypdf import PdfReader


class ParserService:
    def parse_file(self, file_path: str) -> dict[str, Any]:
        path = Path(file_path)
        suffix = path.suffix.lower()

        if suffix in [".txt", ".md", ".markdown"]:
            return self._parse_text(path)
        if suffix == ".csv":
            return self._parse_csv(path)
        if suffix in [".xlsx", ".xlsm", ".xls"]:
            return self._parse_xlsx(path)
        if suffix == ".pdf":
            return self._parse_pdf(path)
        if suffix == ".docx":
            return self._parse_docx(path)
        if suffix == ".pptx":
            return self._parse_pptx(path)

        return {"text": "", "tables": [], "metadata": {"warning": f"Formato não suportado: {suffix}", "filename": path.name, "suffix": suffix}}

    def _parse_text(self, path: Path) -> dict[str, Any]:
        text = path.read_text(encoding="utf-8", errors="ignore")
        return {"text": text, "tables": [], "metadata": {"pages": 1, "format": path.suffix.lower().replace(".", "") or "text", "filename": path.name, "characters": len(text)}}

    def _parse_txt(self, path: Path) -> dict[str, Any]:
        return self._parse_text(path)

    def _parse_csv(self, path: Path) -> dict[str, Any]:
        df = self._read_csv_robust(path)
        full_text = self._df_to_text(df, path.name)
        return {
            "text": full_text,
            "tables": [{
                "sheet": "csv",
                "columns": df.columns.tolist(),
                "row_count": int(df.shape[0]),
                "rows_preview": df.head(20).fillna("").to_dict(orient="records"),
            }],
            "metadata": {"rows": int(df.shape[0]), "columns": int(df.shape[1]), "filename": path.name, "characters": len(full_text)},
        }

    def _parse_xlsx(self, path: Path) -> dict[str, Any]:
        xl = pd.ExcelFile(path)
        texts = []
        tables = []
        total_rows = 0
        for sheet in xl.sheet_names:
            df = xl.parse(sheet, dtype=str).fillna("")
            text = self._df_to_text(df, f"{path.name}:{sheet}")
            texts.append(text)
            total_rows += int(df.shape[0])
            tables.append({"sheet": sheet, "columns": df.columns.tolist(), "row_count": int(df.shape[0]), "rows_preview": df.head(20).to_dict(orient="records")})
        full_text = "\n\n".join(texts)
        return {"text": full_text, "tables": tables, "metadata": {"sheets": xl.sheet_names, "rows": total_rows, "filename": path.name, "characters": len(full_text)}}

    def _parse_pdf(self, path: Path) -> dict[str, Any]:
        reader = PdfReader(str(path))
        pages = []
        for i, page in enumerate(reader.pages, start=1):
            try:
                pages.append(f"[Página {i}]\n{page.extract_text() or ''}")
            except Exception:
                pages.append(f"[Página {i}]\n")
        text = "\n\n".join(pages)
        return {"text": text, "tables": [], "metadata": {"pages": len(reader.pages), "filename": path.name, "characters": len(text)}}

    def _parse_docx(self, path: Path) -> dict[str, Any]:
        doc = Document(str(path))
        texts = [p.text for p in doc.paragraphs if p.text and p.text.strip()]
        # Inclui tabelas DOCX no texto, se houver.
        for table in doc.tables:
            for row in table.rows:
                texts.append(" | ".join(cell.text.strip() for cell in row.cells))
        text = "\n".join(texts)
        return {"text": text, "tables": [], "metadata": {"paragraphs": len(texts), "filename": path.name, "characters": len(text)}}

    def _parse_pptx(self, path: Path) -> dict[str, Any]:
        prs = Presentation(str(path))
        slides = []
        for i, slide in enumerate(prs.slides, start=1):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    texts.append(shape.text)
            slides.append(f"[Slide {i}]\n" + "\n".join(texts))
        text = "\n\n".join(slides)
        return {"text": text, "tables": [], "metadata": {"slides": len(prs.slides), "filename": path.name, "characters": len(text)}}

    def _read_csv_robust(self, path: Path) -> pd.DataFrame:
        encodings = ["utf-8-sig", "utf-8", "latin1", "cp1252"]
        seps = [None, ";", ",", "\t", "|"]
        last_exc: Exception | None = None
        for encoding in encodings:
            for sep in seps:
                try:
                    if sep is None:
                        return pd.read_csv(path, dtype=str, keep_default_na=False, encoding=encoding, sep=None, engine="python")
                    return pd.read_csv(path, dtype=str, keep_default_na=False, encoding=encoding, sep=sep)
                except Exception as exc:
                    last_exc = exc
                    continue
        raise last_exc or RuntimeError("Não foi possível ler o CSV")

    def _df_to_text(self, df: pd.DataFrame, source_name: str) -> str:
        """Converte TODO o dataframe em texto pesquisável.

        A versão anterior usava head(50), o que fazia o Nexus perder ocorrências internas como
        TRANSACAO: dentro de campos longos. Aqui preservamos todas as linhas e todos os valores.
        """
        df = df.fillna("").astype(str)
        columns = [str(c) for c in df.columns.tolist()]
        lines = [f"[Tabela: {source_name}]", f"Colunas: {', '.join(columns)}", f"Total de linhas: {len(df)}", ""]
        for idx, row in df.iterrows():
            parts = [f"{col}: {row.get(col, '')}" for col in columns]
            lines.append(f"[Linha {idx + 1}] " + " | ".join(parts))
        return "\n".join(lines)
